"""Document parser service for extracting text from PDF, DOCX, PPTX files"""

import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)


class DocumentParser:
    """Service for parsing various document formats"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """
        Extract text from PDF file using pdfminer.six
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text content
        """
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(file_path)
            logger.info(f"Extracted {len(text)} characters from PDF")
            return text
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """
        Extract text from DOCX file using python-docx
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Extracted text content
        """
        try:
            from docx import Document
            doc = Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract text from tables
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            tables_text.append(cell.text)
            
            text = '\n'.join(paragraphs + tables_text)
            logger.info(f"Extracted {len(text)} characters from DOCX")
            return text
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            raise ValueError(f"Failed to parse DOCX: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """
        Extract text from PPTX file using python-pptx
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Extracted text content
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            text = '\n'.join(text_parts)
            logger.info(f"Extracted {len(text)} characters from PPTX")
            return text
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            raise ValueError(f"Failed to parse PPTX: {str(e)}")
    
    @staticmethod
    def detect_file_type(file_path: str) -> Optional[str]:
        """
        Detect file type using file extension and magic bytes
        
        Args:
            file_path: Path to file
            
        Returns:
            File type: 'pdf', 'docx', 'pptx', or None
        """
        # Check extension first
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return 'pdf'
        elif ext == '.docx':
            return 'docx'
        elif ext == '.pptx':
            return 'pptx'
        
        # Try magic bytes detection
        try:
            import magic
            mime = magic.from_file(file_path, mime=True)
            
            if mime == 'application/pdf':
                return 'pdf'
            elif mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return 'docx'
            elif mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return 'pptx'
        except Exception as e:
            logger.warning(f"Magic detection failed: {e}")
        
        return None
    
    @classmethod
    def parse_document(cls, file_path: str) -> str:
        """
        Parse document based on file type
        
        Args:
            file_path: Path to document file
            
        Returns:
            Extracted text content
            
        Raises:
            ValueError: If file type is unsupported or parsing fails
        """
        file_type = cls.detect_file_type(file_path)
        
        if file_type == 'pdf':
            return cls.parse_pdf(file_path)
        elif file_type == 'docx':
            return cls.parse_docx(file_path)
        elif file_type == 'pptx':
            return cls.parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type. Please upload PDF, DOCX, or PPTX files.")
